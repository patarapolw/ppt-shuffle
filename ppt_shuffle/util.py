import copy

from pptx.parts.chart import ChartPart
from pptx.parts.embeddedpackage import EmbeddedXlsxPart


def _get_blank_slide_layout(pres):
    layout_items_count = [len(layout.placeholders)
                          for layout in pres.slide_layouts]
    min_items = min(layout_items_count)
    blank_layout_id = layout_items_count.index(min_items)
    return pres.slide_layouts[blank_layout_id]


def duplicate_slide(pres, index):
    source = pres.slides[index]
    blank_slide_layout = _get_blank_slide_layout(pres)
    dest = pres.slides.add_slide(blank_slide_layout)

    for shape in source.shapes:
        newel = copy.deepcopy(shape.element)
        dest.shapes._spTree.insert_element_before(newel, 'p:extLst')

    for key, value in source.part.rels.items():
        # Make sure we don't copy a notesSlide relation as that won't exist
        if "notesSlide" not in value.reltype:
            target = value._target
            # if the relationship was a chart, we need to duplicate the embedded chart part and xlsx
            if "chart" in value.reltype:
                partname = target.package.next_partname(
                    ChartPart.partname_template)
                xlsx_blob = target.chart_workbook.xlsx_part.blob
                target = ChartPart(partname, target.content_type,
                                   copy.deepcopy(target._element), package=target.package)

                target.chart_workbook.xlsx_part = EmbeddedXlsxPart.new(
                    xlsx_blob, target.package)

            dest.part.rels.add_relationship(value.reltype,
                                            target,
                                            value.rId)

    return dest


def delete_slide(presentation, index):
    xml_slides = presentation.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[index])
